"""
Script para criar uma apresentação PowerPoint sobre o problema e a solução do ZF Portal Intelligence Agent
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle=None):
    """Adiciona um slide de título à apresentação"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle
    
    return slide

def add_section_slide(prs, title, content=None):
    """Adiciona um slide de seção à apresentação"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    if content:
        content_shape.text = content
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Adiciona um slide com duas colunas"""
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title_shape = slide.shapes.title
    
    left_shape = slide.placeholders[1]
    right_shape = slide.placeholders[2]
    
    title_shape.text = title
    left_shape.text = left_content
    right_shape.text = right_content
    
    return slide

def add_bullet_slide(prs, title, bullet_points):
    """Adiciona um slide com bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = content_shape.text_frame
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
    
    return slide

def add_table_slide(prs, title, rows, cols, table_data):
    """Adiciona um slide com uma tabela"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Adiciona uma tabela ao slide
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(4.0)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Preenche a tabela com os dados
    for r in range(rows):
        for c in range(cols):
            table.cell(r, c).text = table_data[r][c]
    
    return slide

def create_presentation():
    """Cria a apresentação completa"""
    prs = Presentation()
    
    # Slide 1: Título da apresentação
    add_title_slide(prs, "ZF Portal Intelligence Agent", 
                   "Descrição do Problema e Contextualização")
    
    # Slide 2: O Cenário Atual
    scenario_bullets = [
        "7h30: Chega ao escritório e encontra 50+ emails sobre assuntos urgentes",
        "9h00: Reunião de emergência sobre fluxo de caixa insuficiente",
        "10h30: Descobre que precisará antecipar recebíveis para pagar fornecedores",
        "11h00: Começa a pesquisar soluções de antecipação, ligando para bancos",
        "14h00: Após várias ligações, conseguiu apenas 2 cotações com taxas altas",
        "16h00: Recebe 3 abordagens genéricas por email sobre serviços financeiros",
        "18h30: Encerra o dia com uma solução insatisfatória e caro demais"
    ]
    
    section_slide = add_section_slide(prs, "O Cenário Atual: O Desafio da Prospecção Financeira")
    tf = section_slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "O dia-a-dia do Diretor Financeiro Carlos, de uma rede de varejo com 15 lojas:"
    p.font.bold = True
    
    for bullet in scenario_bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
    
    # Slide 3: Problemas Identificados no Mercado
    problem_slide = add_section_slide(prs, "Problemas Identificados no Mercado")
    tf = problem_slide.placeholders[1].text_frame
    
    # Problema 1
    p = tf.add_paragraph()
    p.text = "1. Processo Manual e Ineficiente"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 70% dos profissionais financeiros gastam +5 horas semanais buscando soluções"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Empresas perdem oportunidades por não encontrar as melhores condições"
    p.level = 1
    
    # Problema 2
    p = tf.add_paragraph()
    p.text = "\n2. Dificuldade de Acesso aos Decisores"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Taxa média de sucesso em contatar decisores financeiros: apenas 12%"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Serviços financeiros enfrentam CAC 5x mais alto que outros setores"
    p.level = 1
    
    # Problema 3
    p = tf.add_paragraph()
    p.text = "\n3. Falta de Personalização"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 82% das abordagens comerciais são genéricas e ignoram contexto específico"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Apenas 3% das comunicações consideram o momento financeiro da empresa"
    p.level = 1
    
    # Slide 4: O Impacto dos Problemas
    impact_slide = add_two_column_slide(prs, "O Impacto dos Problemas", "", "")
    
    # Coluna da esquerda
    left_tf = impact_slide.placeholders[1].text_frame
    p = left_tf.add_paragraph()
    p.text = "Para Empresas que Precisam de Antecipação:"
    p.font.bold = True
    
    bullets1 = [
        "Taxas até 40% mais caras do que o necessário",
        "Tempo médio para encontrar solução: 3-5 dias úteis",
        "Redução de 15-20% na margem de lucro",
        "Atrasos em pagamentos de fornecedores (85% reportam este problema)"
    ]
    
    for bullet in bullets1:
        p = left_tf.add_paragraph()
        p.text = "• " + bullet
    
    # Coluna da direita
    right_tf = impact_slide.placeholders[2].text_frame
    p = right_tf.add_paragraph()
    p.text = "Para Provedores de Serviços Financeiros:"
    p.font.bold = True
    
    bullets2 = [
        "Alto custo de aquisição de clientes",
        "Baixa taxa de conversão (média de 2.7%)",
        "Ineficiência no direcionamento de propostas",
        "Perda de oportunidades de negócio"
    ]
    
    for bullet in bullets2:
        p = right_tf.add_paragraph()
        p.text = "• " + bullet
    
    # Slide 5: Um Exemplo Concreto
    case_study_slide = add_section_slide(prs, "Um Exemplo Concreto: Empresa de Varejo ABC")
    tf = case_study_slide.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Situação:"
    p.font.bold = True
    
    situation_bullets = [
        "Vendas parceladas representam 60% do faturamento",
        "Ciclo médio de recebimento: 45-60 dias",
        "Pagamento a fornecedores: 30 dias"
    ]
    
    for bullet in situation_bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nProblemas Enfrentados:"
    p.font.bold = True
    
    problem_bullets = [
        "Gasta 20+ horas mensais procurando soluções financeiras",
        "Funcionários dedicados apenas para gerenciar antecipações",
        "Usa planilhas manuais para comparar diferentes ofertas",
        "Dificuldade em manter controle sobre taxas negociadas",
        "Média de 4% do faturamento gasto em custos financeiros evitáveis"
    ]
    
    for bullet in problem_bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 1
    
    # Slide 6: Por que este Projeto?
    why_slide = add_section_slide(prs, "Por que este Projeto?")
    tf = why_slide.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Transformação Digital do Processo"
    p.font.bold = True
    
    transformation_bullets = [
        "Potencial de redução de 40% no tempo gasto com prospecção",
        "Economia média de 25% em custos financeiros",
        "Identificação proativa de oportunidades antes de crises de caixa"
    ]
    
    for bullet in transformation_bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nOportunidade de Mercado"
    p.font.bold = True
    
    market_bullets = [
        "R$ 1,8 trilhão/ano em volume de antecipação de recebíveis no Brasil",
        "70% das PMEs buscam soluções para melhorar seu fluxo de caixa",
        "Crescimento de 30% na utilização da antecipação de recebíveis em 2025"
    ]
    
    for bullet in market_bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 1
    
    # Slide 7: A Solução - ZF Portal Intelligence Agent
    solution_slide = add_section_slide(prs, "O ZF Portal Intelligence Agent")
    tf = solution_slide.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Nossa Solução"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "O ZF Portal Intelligence Agent é um sistema de automação inteligente que:"
    
    solution_bullets = [
        "IDENTIFICA decisores financeiros com precisão",
        "PERSONALIZA comunicação baseada no contexto específico",
        "AUTOMATIZA processos de prospecção e follow-up",
        "ANALISA resultados e melhora continuamente"
    ]
    
    for bullet in solution_bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nBenefícios Esperados:"
    p.font.bold = True
    
    benefit_bullets = [
        "Redução do CAC em até 60%",
        "Aumento de taxas de conversão em 300%",
        "Economia de tempo e recursos significativa",
        "Melhores taxas para antecipação de recebíveis"
    ]
    
    for bullet in benefit_bullets:
        p = tf.add_paragraph()
        p.text = "• " + bullet
        p.level = 1
    
    # Slide 8: Tabela resumo do problema e solução
    table_data = [
        ["Problema", "Solução do ZF Portal"],
        ["Processo manual e ineficiente", "Automação inteligente com IA"],
        ["Dificuldade de acesso a decisores", "Identificação precisa via LinkedIn e outras plataformas"],
        ["Falta de personalização", "Comunicação contextualizada para cada perfil"],
        ["Alto custo de aquisição", "Redução significativa do CAC"],
        ["Tempo perdido em prospecção", "Identificação proativa de oportunidades"],
        ["Escolhas financeiras subótimas", "Análise inteligente do mercado"]
    ]
    
    table_slide = add_table_slide(prs, "Resumo do Problema e da Solução", len(table_data), 2, table_data)
    
    # Salvar a apresentação
    prs.save("c:/Users/danie/Desktop/bocao/docs/ZF_Portal_Problema_e_Solucao.pptx")
    return "c:/Users/danie/Desktop/bocao/docs/ZF_Portal_Problema_e_Solucao.pptx"

if __name__ == "__main__":
    create_presentation()
    print("Apresentação criada com sucesso!")
